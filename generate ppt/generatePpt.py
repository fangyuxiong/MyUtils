#coding=utf-8

from __future__ import division
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import glob
import os
import sys

img_path = '/Users/XK/Pictures/www.3dmgame.com/'

ppt_path = '/Users/XK/Desktop/ppt/'

ppt_file = '.pptx'

max_page = 10

max_height = 540

max_width = 720

oneInch = 72

#获取图片的宽和高
def getImageSize(filename):
	img = Image.open(filename)
	return img.size

#将图片顶到头的size
def changeSize(oldSize):
	w = oldSize[0]
	h = oldSize[1]
	nw = 0
	nh = 0
	if (w / h) > (max_width / max_height):
		nw = max_width
		nh = h * max_width / w
	else:
		nh = max_height
		nw = w * max_height / h
	return [nw,nh]

#根据新size获取left 和 top 数组
def getLeftAndTop(newSize):
	w = newSize[0]
	h = newSize[1]
	left = (max_width - w ) / 2
	top = (max_height - h) / 2
	lt = [left,top]
	return lt

def getInches(px):
	return Inches(px / oneInch)

def generatePpt(img_path,filename):
	prs = Presentation()
	index = 0
	for img in glob.glob(img_path + '*'):
		if os.path.isfile(img):
			index = index + 1
			if index > 10:
				break
			size = changeSize(getImageSize(img))
			lt = getLeftAndTop(size)
			blank_slide_layout = prs.slide_layouts[6]
			slide = prs.slides.add_slide(blank_slide_layout)
			pic = slide.shapes.add_picture(img,getInches(lt[0]),getInches(lt[1]),width=getInches(size[0]),height=getInches(size[1]))
	prs.save(ppt_path + filename + ppt_file)
	return 1

def generatePptWithImages(images,filename):
	prs = Presentation()
	index = 0
	for img in images:
		if os.path.isfile(img):
			index = index + 1
			if index > 10:
				break
			size = changeSize(getImageSize(img))
			lt = getLeftAndTop(size)
			blank_slide_layout = prs.slide_layouts[6]
			slide = prs.slides.add_slide(blank_slide_layout)
			pic = slide.shapes.add_picture(img,getInches(lt[0]),getInches(lt[1]),width=getInches(size[0]),height=getInches(size[1]))
	prs.save(ppt_path + filename + ppt_file)
	return 1

if __name__ == '__main__':
	l = len(sys.argv)
	if l == 2:
		generatePpt(img_path,sys.argv[1])
	elif l == 3:
		generatePpt(sys.argv[1],sys.argv[2])


